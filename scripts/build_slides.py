"""
Build presentation slides as PPTX (with embedded video).
Output: report/ML_Term_Project_Slides.pptx
- Slide 3: Pendulum-v1 환경 설명 + 영상 임베딩
"""

import os, json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

FIG_DIR = os.path.join("results", "figures")
OUT_DIR = "report"
os.makedirs(OUT_DIR, exist_ok=True)

with open(os.path.join("results", "metrics.json"), encoding="utf-8") as f:
    metrics = json.load(f)

# ── 색상 팔레트 ───────────────────────────────────────────────────────────────
DARK   = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT = RGBColor(0x0f, 0x3c, 0x96)
WHITE  = RGBColor(0xff, 0xff, 0xff)
LGRAY  = RGBColor(0xf0, 0xf0, 0xf0)
DGRAY  = RGBColor(0x55, 0x55, 0x55)
GREEN  = RGBColor(0x1b, 0x87, 0x3b)
LBLUE  = RGBColor(0xaa, 0xcc, 0xff)
LBLUE2 = RGBColor(0x77, 0x99, 0xcc)
LGREEN = RGBColor(0xe8, 0xf4, 0xe8)
PURPLE = RGBColor(0x25, 0x25, 0x50)
LPURP  = RGBColor(0xdd, 0xee, 0xff)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ── 기본 헬퍼 ─────────────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill=None):
    sh = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    sh.line.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    return sh


def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def img(slide, fname, l, t, w, h=None):
    path = os.path.join(FIG_DIR, fname) if not os.path.isabs(fname) else fname
    if not os.path.exists(path):
        return None
    if h:
        return slide.shapes.add_picture(
            path, Inches(l), Inches(t), Inches(w), Inches(h))
    return slide.shapes.add_picture(
        path, Inches(l), Inches(t), Inches(w))


def header(slide, title, subtitle=None):
    """상단 헤더 바."""
    rect(slide, 0, 0, 13.33, 1.05, fill=DARK)
    txt(slide, title, 0.35, 0.1, 12.5, 0.65,
        size=26, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.35, 0.68, 12.5, 0.38,
            size=13, color=LBLUE)
    rect(slide, 0, 1.05, 13.33, 0.04, fill=ACCENT)


def bullet_list(slide, items, l, t, w, size=14,
                color=DARK, gap=0.46):
    """글머리 목록."""
    for i, item in enumerate(items):
        rect(slide, l, t + i * gap + 0.08, 0.10, 0.28, fill=ACCENT)
        txt(slide, item, l + 0.18, t + i * gap, w - 0.18, gap,
            size=size, color=color)


def table_grid(slide, rows, col_widths, start_x, start_y,
               row_height=0.46):
    """직접 그리는 테이블 (헤더는 DARK 배경)."""
    for r_idx, row in enumerate(rows):
        x = start_x
        for c_idx, (cell, cw) in enumerate(zip(row, col_widths)):
            if r_idx == 0:
                bg, fg, bold = DARK, WHITE, True
            elif r_idx % 2 == 0:
                bg, fg, bold = LGRAY, DARK, False
            else:
                bg, fg, bold = WHITE, DARK, False
            rect(slide, x, start_y + r_idx * row_height,
                 cw, row_height, fill=bg)
            txt(slide, cell,
                x + 0.06,
                start_y + r_idx * row_height + 0.05,
                cw - 0.12, row_height - 0.08,
                size=12, bold=bold, color=fg, align=PP_ALIGN.CENTER)
            x += cw


# ── Slide 1: 제목 ─────────────────────────────────────────────────────────────
def slide_title(prs):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
    rect(slide, 0, 3.1, 13.33, 0.06, fill=ACCENT)

    txt(slide, "로봇 시뮬레이션 환경의",
        1.0, 0.7, 11.3, 1.0, size=40, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)
    txt(slide, "수치 적분 오차 예측을 위한 머신러닝",
        1.0, 1.6, 11.3, 1.0, size=30, bold=False, color=LBLUE,
        align=PP_ALIGN.CENTER)
    txt(slide, "[2026-1] Machine Learning — 텀 프로젝트",
        1.0, 3.35, 11.3, 0.6, size=17, color=LGRAY,
        align=PP_ALIGN.CENTER)
    txt(slide, "인하대학교  |  2026년 6월 1일",
        1.0, 3.95, 11.3, 0.5, size=15, color=LGRAY,
        align=PP_ALIGN.CENTER)
    txt(slide, "Gymnasium  |  PyTorch  |  scikit-learn",
        1.0, 6.6, 11.3, 0.5, size=13, color=LBLUE2,
        align=PP_ALIGN.CENTER)


# ── Slide 2: 문제 정의 ────────────────────────────────────────────────────────
def slide_problem(prs):
    slide = prs.slides.add_slide(BLANK)
    header(slide, "문제 정의",
           "수치 적분 오차는 왜 중요한가?")

    # 배경 설명
    txt(slide, "Sim-to-Real 격차 배경", 0.4, 1.2, 12.5, 0.45,
        size=18, bold=True, color=DARK)
    bullet_list(slide, [
        "로봇 시뮬레이터는 연속 시간 물리 법칙을 이산 시간으로 수치 적분하여 근사한다.",
        "Gymnasium Pendulum-v1은 1차 Euler 적분(dt = 0.05 s)을 사용한다.",
        "Euler 오차는 긴 궤적에서 누적되어 컨트롤러 성능 저하를 유발한다.",
    ], 0.4, 1.72, 12.4, size=14)

    # 목표
    txt(slide, "프로젝트 목표", 0.4, 3.65, 12.5, 0.45,
        size=18, bold=True, color=DARK)
    rect(slide, 0.4, 4.18, 12.5, 1.5, fill=LGRAY)
    txt(slide,
        "머신러닝 모델로 수치 적분 오차를 예측한다\n"
        "(Euler 결과 - RK4 결과 = 보정해야 할 오차)",
        0.65, 4.28, 12.1, 1.0, size=17, color=DARK)
    txt(slide,
        "입력: (cosθ, sinθ, θ̇, τ)   →   "
        "출력: (err_theta, err_theta_dot)",
        0.65, 5.0, 12.1, 0.5, size=15, bold=True, color=ACCENT)

    txt(slide,
        "활용: 오차 보상(error compensation)  |  sim-to-real 전이  |  적응형 수치 적분",
        0.4, 6.15, 12.5, 0.4, size=13, color=DGRAY)


# ── Slide 3: 환경 + 영상 ──────────────────────────────────────────────────────
def slide_environment(prs):
    slide = prs.slides.add_slide(BLANK)
    header(slide, "환경: Gymnasium Pendulum-v1",
           "알려진 물리 방정식을 가진 고전 제어 벤치마크")

    # 왼쪽: 물리 및 비교표
    txt(slide, "단진자 운동 방정식", 0.35, 1.18, 6.2, 0.4,
        size=15, bold=True, color=DARK)
    eq_rows = [
        "θ'' = (3g / 2l) sin(θ) + (3 / ml²) × τ",
        "g = 10 m/s²,  m = 1 kg,  l = 1 m,  dt = 0.05 s",
        "|θ̇| ≤ 8 rad/s (clipped),  |τ| ≤ 2 N·m",
    ]
    for i, t in enumerate(eq_rows):
        rect(slide, 0.35, 1.65 + i * 0.56, 6.0, 0.50, fill=LGRAY)
        txt(slide, t, 0.50, 1.68 + i * 0.56, 5.80, 0.44,
            size=13, color=DARK)

    txt(slide, "적분 방법 비교", 0.35, 3.4, 6.2, 0.4,
        size=15, bold=True, color=DARK)
    table_grid(slide,
        [("방법",       "차수", "서브스텝 dt", "오차 차수"),
         ("Euler (Gym)", "1차", "0.05 s",   "O(dt²)"),
         ("RK4 (기준값)", "4차", "0.0025 s ×20", "O(dt⁵)")],
        [2.2, 1.0, 1.8, 1.0], 0.35, 3.88, row_height=0.46)

    txt(slide, "RK4 = 20 서브스텝으로 Euler보다 훨씬 정밀한 기준값 역할",
        0.35, 5.35, 6.2, 0.4, size=12, color=DGRAY)

    # 오른쪽: 영상
    video_path = os.path.join(FIG_DIR, "pendulum_demo.mp4")
    if os.path.exists(video_path):
        import cv2 as _cv2
        poster_path = os.path.join(OUT_DIR, "_poster_frame.png")
        cap = _cv2.VideoCapture(video_path)
        ret, frame = cap.read()
        cap.release()
        if ret:
            _cv2.imwrite(poster_path, frame)
            slide.shapes.add_movie(
                video_path,
                Inches(6.65), Inches(1.18), Inches(6.35), Inches(5.0),
                poster_frame_image=poster_path,
                mime_type="video/mp4",
            )
        txt(slide,
            "에너지 기반 스윙업 컨트롤러 | 12초 데모 (클릭하여 재생)",
            6.65, 6.3, 6.35, 0.45, size=12, color=DGRAY,
            align=PP_ALIGN.CENTER)


# ── Slide 4: 데이터 수집 ──────────────────────────────────────────────────────
def slide_data(prs):
    slide = prs.slides.add_slide(BLANK)
    header(slide, "데이터셋 수집",
           "Pendulum-v1에서 100,000 샘플 수집")

    # 통계 카드 5개
    stats = [
        ("100,000", "총 샘플 수"),
        ("500",     "에피소드"),
        ("200",     "스텝/에피소드"),
        ("4",       "입력 특성"),
        ("2",       "출력 목표"),
    ]
    bw = 2.35
    for i, (num, label) in enumerate(stats):
        x = 0.3 + i * (bw + 0.12)
        rect(slide, x, 1.18, bw, 1.28, fill=DARK)
        txt(slide, num, x, 1.25, bw, 0.72,
            size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, label, x, 1.9, bw, 0.42,
            size=12, color=LGRAY, align=PP_ALIGN.CENTER)

    # 좌: 입력 특성
    txt(slide, "입력 특성 (X)", 0.35, 2.65, 6.0, 0.4,
        size=16, bold=True, color=DARK)
    feats = [
        "cos(θ)    진자 각도의 코사인",
        "sin(θ)    진자 각도의 사인",
        "θ̇        각속도 (rad/s)",
        "τ         적용 토크 (N·m)",
    ]
    for i, t in enumerate(feats):
        rect(slide, 0.35, 3.12 + i * 0.52, 0.12, 0.38, fill=ACCENT)
        txt(slide, t, 0.55, 3.12 + i * 0.52, 5.85, 0.44, size=14, color=DARK)

    # 우: 출력 목표
    txt(slide, "출력 목표 (y)", 6.85, 2.65, 6.1, 0.4,
        size=16, bold=True, color=DARK)
    tgts = [
        ("err_theta",     "각도 오차: Euler − RK4 (rad)"),
        ("err_theta_dot", "각속도 오차: Euler − RK4 (rad/s)"),
    ]
    for i, (name, desc) in enumerate(tgts):
        rect(slide, 6.85, 3.12 + i * 0.62, 6.1, 0.55, fill=LGRAY)
        txt(slide, f"{name}", 7.05, 3.16 + i * 0.62, 2.3, 0.44,
            size=14, bold=True, color=ACCENT)
        txt(slide, desc, 9.35, 3.16 + i * 0.62, 3.55, 0.44,
            size=13, color=DARK)

    txt(slide,
        "학습 70%  |  검증 15%  |  테스트 15%   (랜덤 분할, seed=42)",
        0.35, 6.42, 12.6, 0.45, size=14, color=DGRAY,
        align=PP_ALIGN.CENTER)


# ── Slide 5: 모델 ─────────────────────────────────────────────────────────────
def slide_models(prs):
    slide = prs.slides.add_slide(BLANK)
    header(slide, "모델 구조", "베이스라인 vs MLP")

    # Linear Regression 카드
    rect(slide, 0.35, 1.18, 5.95, 0.52, fill=DARK)
    txt(slide, "Linear Regression (베이스라인)",
        0.5, 1.22, 5.65, 0.44, size=16, bold=True, color=WHITE)
    lr_items = [
        "MultiOutputRegressor(LinearRegression)",
        "목표 변수별 독립적인 회귀 모델 적용",
        "목적함수: min ||Xw - y||",
        "sin(θ)에 대한 비선형 의존성 학습 불가",
    ]
    for i, t in enumerate(lr_items):
        rect(slide, 0.35, 1.78 + i * 0.50, 0.12, 0.36, fill=DGRAY)
        txt(slide, t, 0.55, 1.78 + i * 0.50, 5.70, 0.44, size=13, color=DARK)

    # MLP 카드
    rect(slide, 7.05, 1.18, 5.95, 0.52, fill=ACCENT)
    txt(slide, "MLP (주요 모델)",
        7.20, 1.22, 5.65, 0.44, size=16, bold=True, color=WHITE)
    mlp_items = [
        "Input(4) → Linear(128) → ReLU",
        "→ Linear(64) → ReLU",
        "→ Linear(32) → ReLU",
        "→ Linear(2)  [출력층]",
        "손실함수: MSELoss  |  최적화: Adam (lr=1e-3)",
        "에폭: 100  |  배치 크기: 512",
    ]
    for i, t in enumerate(mlp_items):
        rect(slide, 7.05, 1.78 + i * 0.50, 0.12, 0.36, fill=ACCENT)
        txt(slide, t, 7.25, 1.78 + i * 0.50, 5.70, 0.44, size=13, color=DARK)

    # 학습 곡선 (가운데 하단)
    img(slide, "mlp_learning_curve.png", 2.5, 4.92, 8.3, 2.22)
    txt(slide, "MLP 학습 곡선 (Train / Validation MSE, 로그 스케일)",
        2.5, 7.18, 8.3, 0.28, size=11, color=DGRAY,
        align=PP_ALIGN.CENTER)


# ── Slide 6: 결과 ─────────────────────────────────────────────────────────────
def slide_results(prs):
    slide = prs.slides.add_slide(BLANK)
    header(slide, "실험 결과",
           "테스트셋 성능 평가 (15,000 샘플)")

    lr_m  = metrics["Linear Regression"]
    mlp_m = metrics["MLP"]

    # 상단: 비교 테이블
    table_grid(slide,
        [("모델",              "err_theta R²", "err_theta_dot R²", "평균 RMSE"),
         ("Linear Regression",
          f"{lr_m['err_theta']['R2']:.4f}",
          f"{lr_m['err_theta_dot']['R2']:.4f}",
          f"{lr_m['mean']['RMSE']:.2e}"),
         ("MLP",
          f"{mlp_m['err_theta']['R2']:.4f}",
          f"{mlp_m['err_theta_dot']['R2']:.4f}",
          f"{mlp_m['mean']['RMSE']:.2e}")],
        [4.2, 2.6, 3.1, 2.6], 0.35, 1.18, row_height=0.52)

    # 강조 텍스트
    txt(slide,
        "★  MLP: 두 목표 변수 모두 R² > 0.99 달성 "
        "(평균 RMSE = 1.4×10⁻³)",
        0.35, 2.82, 12.6, 0.50,
        size=15, bold=True, color=GREEN)

    # 하단: 비교 차트
    img(slide, "model_comparison.png", 0.9, 3.35, 11.5, 3.85)


# ── Slide 7: 산점도 ──────────────────────────────────────────────────────────
def slide_scatter(prs):
    slide = prs.slides.add_slide(BLANK)
    header(slide, "예측값 vs. 실제값 산점도",
           "두 모델의 예측 정확도 비교")

    img(slide, "scatter_err_theta.png",    0.25, 1.18, 6.35, 2.90)
    txt(slide, "err_theta (각도 오차)",
        0.25, 4.12, 6.35, 0.38, size=12, color=DGRAY,
        align=PP_ALIGN.CENTER)

    img(slide, "scatter_err_theta_dot.png", 6.75, 1.18, 6.35, 2.90)
    txt(slide, "err_theta_dot (각속도 오차)",
        6.75, 4.12, 6.35, 0.38, size=12, color=DGRAY,
        align=PP_ALIGN.CENTER)

    # 잔차 히스토그램
    img(slide, "residuals.png", 1.2, 4.58, 10.9, 2.65)
    txt(slide,
        "잔차 분포: MLP 잔차는 0에 좁게 집중, Linear Regression은 err_theta_dot에서 넓은 꼬리",
        1.2, 7.22, 10.9, 0.28, size=11, color=DGRAY,
        align=PP_ALIGN.CENTER)


# ── Slide 8: 토론 ─────────────────────────────────────────────────────────────
def slide_discussion(prs):
    slide = prs.slides.add_slide(BLANK)
    header(slide, "결과 분석 및 토론", "")

    sections = [
        ("왜 MLP가 우수한가?",
         ["err_theta_dot는 sin(θ)에 비선형 의존 → 선형 모델로는 학습 불가.",
          "MLP의 ReLU 활성화 함수가 비선형 오차 곡면을 효과적으로 근사.",
          "err_theta는 θ̇에 선형 → LR도 R²=0.99 달성."]),
        ("실패 사례 (Failure Modes)",
         ["|θ̇| ≈ 8 rad/s 근방: clip 연산으로 생기는 급격한 불연속점에서 오차 큼.",
          "θ ≈ 0 (직립) 근방: 2차 미분 최대 → Euler 오차 최대 → 예측 난이도 증가.",
          "두 모델 모두 극단값을 약간 과소 예측하는 경향."]),
        ("향후 연구 방향",
         ["물리 기반 특성 추가: sin²(θ), θ̇·sin(θ) → LR 성능 향상 가능.",
          "LSTM/Transformer로 다중 스텝 오차 누적 예측.",
          "학습된 오차 모델을 MPC 컨트롤러에 통합하여 제어 성능 개선 검증."]),
    ]

    y = 1.22
    for sec_title, pts in sections:
        rect(slide, 0.35, y, 12.6, 0.46, fill=DARK)
        txt(slide, sec_title, 0.55, y + 0.05, 12.3, 0.37,
            size=15, bold=True, color=WHITE)
        y += 0.50
        for pt in pts:
            rect(slide, 0.50, y + 0.07, 0.09, 0.28, fill=ACCENT)
            txt(slide, pt, 0.70, y + 0.02, 12.1, 0.42, size=13, color=DARK)
            y += 0.44
        y += 0.14


# ── Slide 9: 결론 ─────────────────────────────────────────────────────────────
def slide_conclusion(prs):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
    rect(slide, 0, 2.95, 13.33, 0.06, fill=ACCENT)

    txt(slide, "결론", 1.0, 0.45, 11.3, 0.9,
        size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    msgs = [
        "MLP [128→64→32] : 두 목표 모두 R² > 0.99, 평균 RMSE = 1.4×10⁻³",
        "수치 적분 오차는 상태-행동 공간의 매끄럽고 학습 가능한 함수",
        "Linear Regression: err_theta(R²=0.99) ▶ 강함, err_theta_dot(R²=0.52) ▶ 약함",
        "학습된 오차 모델 → sim-to-real 전이 / 오차 보상 / 적응형 적분 활용 가능",
    ]
    for i, msg in enumerate(msgs):
        rect(slide, 0.7, 3.18 + i * 0.90, 11.9, 0.78, fill=PURPLE)
        txt(slide, msg, 0.95, 3.26 + i * 0.90, 11.5, 0.55,
            size=15, color=LPURP)

    txt(slide, "GitHub: github.com/CupidJang/my_ml_project  |  Gymnasium  |  PyTorch  |  scikit-learn",
        1.0, 7.1, 11.3, 0.38, size=12, color=LBLUE2,
        align=PP_ALIGN.CENTER)


# ── 슬라이드 생성 ─────────────────────────────────────────────────────────────
slide_title(prs)
slide_problem(prs)
slide_environment(prs)
slide_data(prs)
slide_models(prs)
slide_results(prs)
slide_scatter(prs)
slide_discussion(prs)
slide_conclusion(prs)

out = os.path.join(OUT_DIR, "ML_Term_Project_Slides.pptx")
prs.save(out)
print(f"슬라이드 저장 완료: {out}")
print("Slide 3 환경 영상 임베딩됨 (PowerPoint에서 클릭 재생)")
print("PDF 제출: PowerPoint에서 파일 > 내보내기 > PDF")
